import os
import glob
import base64
import tempfile

from langchain_core.documents import Document as LC_Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings, OllamaLLM
from langchain_core.messages import HumanMessage
from langchain_community.retrievers import BM25Retriever
import pickle
from rich.console import Console

console = Console()

APP_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.abspath(os.path.join(APP_DIR, 'data'))
RAW_DOCS_DIR = os.path.join(DATA_DIR, 'raw_docs')
CHROMA_DB_DIR = os.path.join(DATA_DIR, 'chroma_db')
BM25_INDEX_PATH = os.path.join(DATA_DIR, 'bm25_index.pkl')

def describe_image(image_path):
    """Uses the local Vision model to describe an extracted image or flowchart."""
    try:
        llm = OllamaLLM(model="llama3.2-vision:11b", num_thread=8)
        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        
        message = HumanMessage(
            content=[
                {"type": "text", "text": "Describe this flowchart, diagram, or image in detail. Extract any text visible and explain the logical steps or data shown. Prioritize accuracy and technical detail. Do not use markdown blocks."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded_string}"}}
            ]
        )
        return llm.invoke([message])
    except Exception as e:
        console.print(f"[dim red]Warning: Vision model failed on {os.path.basename(image_path)}: {e}[/dim red]")
        return "Image description unavailable."

def extract_docx_with_vision(file_path):
    """Parses Word docs, extracts text, and sends any embedded images to the Vision LLM."""
    import docx2txt
    
    docs = []
    with tempfile.TemporaryDirectory() as temp_dir:
        text = docx2txt.process(file_path, temp_dir)
        images = glob.glob(os.path.join(temp_dir, "*"))
        
        append_text = ""
        for img in images:
            if img.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                console.print(f"  [dim]Analyzing DOCX image/flowchart with Vision Model...[/dim]")
                desc = describe_image(img)
                append_text += f"\n\n[Embedded Image/Diagram Description: {desc}]\n"
                
        full_text = text + append_text
        if full_text.strip():
            docs.append(LC_Document(page_content=full_text, metadata={"source": file_path}))
    return docs

def extract_pptx_with_vision(file_path):
    """Parses PPTX slides, extracts text, and sends any embedded images to the Vision LLM."""
    from pptx import Presentation
    
    docs = []
    prs = Presentation(file_path)
    with tempfile.TemporaryDirectory() as temp_dir:
        for i, slide in enumerate(prs.slides):
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    img_path = os.path.join(temp_dir, f"slide_{i+1}_{image.filename}")
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    console.print(f"  [dim]Analyzing Slide {i+1} image/flowchart with Vision Model...[/dim]")
                    desc = describe_image(img_path)
                    slide_text.append(f"\n[Embedded Image/Diagram Description: {desc}]")
                    
            content = "\n".join(slide_text).strip()
            if content:
                docs.append(LC_Document(page_content=content, metadata={"source": file_path, "slide": i+1}))
    return docs

def extract_excel(file_path):
    """Parses Excel files and converts sheets into markdown table format for semantic search."""
    import pandas as pd
    
    docs = []
    try:
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            # Convert tabular data into semantic text format
            csv_text = df.to_markdown(index=False)
            content = f"Excel Data - Sheet: {sheet_name}\n\n{csv_text}"
            docs.append(LC_Document(page_content=content, metadata={"source": file_path, "sheet": sheet_name}))
    except Exception as e:
        console.print(f"[bold red]Failed to process Excel file {os.path.basename(file_path)}: {e}[/bold red]")
    return docs

def extract_pdf_with_vision(file_path):
    """Parses PDF docs, extracts text, and sends any embedded images to the Vision LLM."""
    import fitz  # PyMuPDF
    
    docs = []
    try:
        doc = fitz.open(file_path)
        with tempfile.TemporaryDirectory() as temp_dir:
            for i, page in enumerate(doc):
                text = page.get_text()
                
                image_list = page.get_images(full=True)
                append_text = ""
                for image_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    if image_ext.lower() in ['png', 'jpg', 'jpeg']:
                        img_path = os.path.join(temp_dir, f"page_{i+1}_img_{image_index}.{image_ext}")
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)
                        
                        console.print(f"  [dim]Analyzing PDF Page {i+1} image/diagram with Vision Model...[/dim]")
                        desc = describe_image(img_path)
                        append_text += f"\n\n[Embedded Image/Diagram Description: {desc}]\n"
                    
                full_text = text + append_text
                if full_text.strip():
                    docs.append(LC_Document(page_content=full_text, metadata={"source": file_path, "page": i+1}))
    except Exception as e:
        console.print(f"[bold red]Failed to process PDF {os.path.basename(file_path)}: {e}[/bold red]")
    return docs

def load_documents():
    docs = []
    os.makedirs(RAW_DOCS_DIR, exist_ok=True)
    
    word_files = glob.glob(os.path.join(RAW_DOCS_DIR, "*.docx"))
    ppt_files = glob.glob(os.path.join(RAW_DOCS_DIR, "*.pptx"))
    excel_files = glob.glob(os.path.join(RAW_DOCS_DIR, "*.xlsx")) + glob.glob(os.path.join(RAW_DOCS_DIR, "*.xls"))
    pdf_files = glob.glob(os.path.join(RAW_DOCS_DIR, "*.pdf"))
    
    if not any([word_files, ppt_files, excel_files, pdf_files]):
        console.print("[yellow]No supported files (.docx, .pptx, .xlsx, .pdf) found in data/raw_docs/[/yellow]")
        return docs
        
    console.print(f"[cyan]Found {len(pdf_files)} PDFs, {len(word_files)} Word docs, {len(ppt_files)} PowerPoints, and {len(excel_files)} Excel files.[/cyan]")
    
    for file_path in word_files:
        console.print(f"[bold white]Processing {os.path.basename(file_path)}...[/bold white]")
        docs.extend(extract_docx_with_vision(file_path))
        
    for file_path in ppt_files:
        console.print(f"[bold white]Processing {os.path.basename(file_path)}...[/bold white]")
        docs.extend(extract_pptx_with_vision(file_path))
        
        
    for file_path in excel_files:
        console.print(f"[bold white]Processing {os.path.basename(file_path)}...[/bold white]")
        docs.extend(extract_excel(file_path))
        
    for file_path in pdf_files:
        console.print(f"[bold white]Processing {os.path.basename(file_path)}...[/bold white]")
        docs.extend(extract_pdf_with_vision(file_path))
        
    return docs

def split_documents(documents):
    if not documents:
        return []
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=250,
        length_function=len
    )
    chunks = text_splitter.split_documents(documents)
    console.print(f"[cyan]Split {len(documents)} document sections into {len(chunks)} searchable chunks.[/cyan]")
    return chunks

def ingest():
    console.print("\n[bold cyan]─── IDE Expert Agent: Advanced Knowledge Ingestion ───[/bold cyan]")
    
    documents = load_documents()
    if not documents:
        console.print("[bold red]Ingestion aborted: No valid content recovered.[/bold red]")
        return
        
    chunks = split_documents(documents)
    
    console.print("\n[dim]Initializing Embedding Model and ChromaDB...[/dim]")
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    
    console.print(f"[dim]Writing chunks to Vector Database... This may take a few minutes.[/dim]")
    vector_db = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=CHROMA_DB_DIR
    )
    
    console.print("\n[dim]Building BM25 Keyword Search Index...[/dim]")
    bm25_retriever = BM25Retriever.from_documents(chunks)
    with open(BM25_INDEX_PATH, 'wb') as f:
        pickle.dump(bm25_retriever, f)
    
    console.print(f"\n[bold green]Success! Advanced Hybrid Knowledge base updated across {len(chunks)} chunks.[/bold green]")
    console.print(f"Vector Database stored at: {CHROMA_DB_DIR}")
    console.print(f"Keyword Index stored at: {BM25_INDEX_PATH}")

if __name__ == "__main__":
    ingest()
