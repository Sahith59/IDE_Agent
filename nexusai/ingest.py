"""
RAG ingestion pipeline — indexes documents into ~/.nexus/chroma_db and bm25_index.pkl
Usage: nexus-ingest [/path/to/docs/dir]
       Defaults to ~/.nexus/raw_docs/ if no path given.
"""
import os
import sys
import glob
import base64
import tempfile

from nexusai import config as _cfg
_cfg.apply_env()

from langchain_core.documents import Document as LC_Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings, OllamaLLM
from langchain_core.messages import HumanMessage
from langchain_community.retrievers import BM25Retriever
import pickle
from rich.console import Console

console = Console()


def describe_image(image_path):
    try:
        llm = OllamaLLM(model="llama3.2-vision:11b", num_thread=8)
        with open(image_path, "rb") as image_file:
            encoded = base64.b64encode(image_file.read()).decode("utf-8")
        message = HumanMessage(content=[
            {"type": "text", "text": "Describe this flowchart, diagram, or image in detail. Extract any text visible and explain the logical steps or data shown. Prioritize accuracy and technical detail. Do not use markdown blocks."},
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded}"}}
        ])
        return llm.invoke([message])
    except Exception as e:
        console.print(f"[dim red]Warning: Vision model failed on {os.path.basename(image_path)}: {e}[/dim red]")
        return "Image description unavailable."


def extract_docx_with_vision(file_path):
    import docx2txt
    docs = []
    with tempfile.TemporaryDirectory() as tmp:
        text   = docx2txt.process(file_path, tmp)
        images = glob.glob(os.path.join(tmp, "*"))
        extra  = ""
        for img in images:
            if img.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                console.print(f"  [dim]Analyzing DOCX image with Vision Model...[/dim]")
                extra += f"\n\n[Embedded Image Description: {describe_image(img)}]\n"
        full = text + extra
        if full.strip():
            docs.append(LC_Document(page_content=full, metadata={"source": file_path}))
    return docs


def extract_pptx_with_vision(file_path):
    from pptx import Presentation
    docs = []
    prs  = Presentation(file_path)
    with tempfile.TemporaryDirectory() as tmp:
        for i, slide in enumerate(prs.slides):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    img_path = os.path.join(tmp, f"slide_{i+1}_{image.filename}")
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    console.print(f"  [dim]Analyzing Slide {i+1} image with Vision Model...[/dim]")
                    texts.append(f"\n[Embedded Image Description: {describe_image(img_path)}]")
            content = "\n".join(texts).strip()
            if content:
                docs.append(LC_Document(page_content=content,
                                        metadata={"source": file_path, "slide": i+1}))
    return docs


def extract_excel(file_path):
    import pandas as pd
    docs = []
    try:
        xls = pd.ExcelFile(file_path)
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet)
            if df.empty:
                continue
            docs.append(LC_Document(
                page_content=f"Excel Data - Sheet: {sheet}\n\n{df.to_markdown(index=False)}",
                metadata={"source": file_path, "sheet": sheet}
            ))
    except Exception as e:
        console.print(f"[bold red]Failed to process Excel {os.path.basename(file_path)}: {e}[/bold red]")
    return docs


def extract_pdf_with_vision(file_path):
    import fitz
    docs = []
    try:
        doc = fitz.open(file_path)
        with tempfile.TemporaryDirectory() as tmp:
            for i, page in enumerate(doc):
                text  = page.get_text()
                extra = ""
                for img_idx, img in enumerate(page.get_images(full=True)):
                    base_image = doc.extract_image(img[0])
                    ext        = base_image["ext"].lower()
                    if ext in ("png", "jpg", "jpeg"):
                        img_path = os.path.join(tmp, f"page_{i+1}_img_{img_idx}.{ext}")
                        with open(img_path, "wb") as f:
                            f.write(base_image["image"])
                        console.print(f"  [dim]Analyzing PDF Page {i+1} image with Vision Model...[/dim]")
                        extra += f"\n\n[Embedded Image Description: {describe_image(img_path)}]\n"
                full = text + extra
                if full.strip():
                    docs.append(LC_Document(page_content=full,
                                            metadata={"source": file_path, "page": i+1}))
    except Exception as e:
        console.print(f"[bold red]Failed to process PDF {os.path.basename(file_path)}: {e}[/bold red]")
    return docs


def extract_markdown(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        if content.strip():
            return [LC_Document(page_content=content,
                                metadata={"source": file_path, "type": "markdown"})]
    except Exception as e:
        console.print(f"[bold red]Failed to process Markdown {os.path.basename(file_path)}: {e}[/bold red]")
    return []


def load_documents(raw_docs_dir: str):
    docs = []
    os.makedirs(raw_docs_dir, exist_ok=True)

    word_files  = glob.glob(os.path.join(raw_docs_dir, "*.docx"))
    ppt_files   = glob.glob(os.path.join(raw_docs_dir, "*.pptx"))
    excel_files = glob.glob(os.path.join(raw_docs_dir, "*.xlsx")) + \
                  glob.glob(os.path.join(raw_docs_dir, "*.xls"))
    pdf_files   = glob.glob(os.path.join(raw_docs_dir, "*.pdf"))
    md_files    = glob.glob(os.path.join(raw_docs_dir, "*.md"))

    if not any([word_files, ppt_files, excel_files, pdf_files, md_files]):
        console.print(f"[yellow]No supported files found in {raw_docs_dir}[/yellow]")
        console.print("[dim]Supported: .docx .pptx .xlsx .pdf .md[/dim]")
        return docs

    console.print(
        f"[cyan]Found {len(pdf_files)} PDFs, {len(word_files)} Word, "
        f"{len(ppt_files)} PowerPoint, {len(excel_files)} Excel, "
        f"{len(md_files)} Markdown.[/cyan]"
    )

    for fp in word_files:
        console.print(f"[bold white]Processing {os.path.basename(fp)}...[/bold white]")
        docs.extend(extract_docx_with_vision(fp))
    for fp in ppt_files:
        console.print(f"[bold white]Processing {os.path.basename(fp)}...[/bold white]")
        docs.extend(extract_pptx_with_vision(fp))
    for fp in excel_files:
        console.print(f"[bold white]Processing {os.path.basename(fp)}...[/bold white]")
        docs.extend(extract_excel(fp))
    for fp in pdf_files:
        console.print(f"[bold white]Processing {os.path.basename(fp)}...[/bold white]")
        docs.extend(extract_pdf_with_vision(fp))
    for fp in md_files:
        console.print(f"[bold white]Processing {os.path.basename(fp)}...[/bold white]")
        docs.extend(extract_markdown(fp))

    return docs


def ingest(docs_dir: str | None = None):
    raw_docs_dir = docs_dir or _cfg.RAW_DOCS_DIR
    console.print(f"\n[bold cyan]─── Nexus: Knowledge Ingestion ───[/bold cyan]")
    console.print(f"[dim]Source: {raw_docs_dir}[/dim]")

    documents = load_documents(raw_docs_dir)
    if not documents:
        console.print("[bold red]Ingestion aborted: No valid content found.[/bold red]")
        return

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500, chunk_overlap=250, length_function=len
    )
    chunks = splitter.split_documents(documents)
    console.print(f"[cyan]Split into {len(chunks)} searchable chunks.[/cyan]")

    console.print("\n[dim]Initializing ChromaDB...[/dim]")
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=_cfg.CHROMA_DB_DIR,
    )

    console.print("\n[dim]Building BM25 keyword index...[/dim]")
    bm25 = BM25Retriever.from_documents(chunks)
    with open(_cfg.BM25_INDEX_PATH, "wb") as f:
        pickle.dump(bm25, f)

    console.print(f"\n[bold green]Done! {len(chunks)} chunks indexed.[/bold green]")
    console.print(f"Vector DB: {_cfg.CHROMA_DB_DIR}")
    console.print(f"BM25 index: {_cfg.BM25_INDEX_PATH}")


def main():
    docs_dir = sys.argv[1] if len(sys.argv) > 1 else None
    ingest(docs_dir)


if __name__ == "__main__":
    main()
